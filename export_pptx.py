"""
export_pptx.py — Generates presentation.pptx from the slide data.
Run: python export_pptx.py
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours (matching cream/light theme) ──────────────────
C_BG      = RGBColor(0xF0, 0xEB, 0xE3)   # cream bg
C_DARK    = RGBColor(0x1C, 0x1A, 0x16)   # charcoal
C_MID     = RGBColor(0x4A, 0x40, 0x35)   # body text brown
C_MUTED   = RGBColor(0x8A, 0x7D, 0x6E)   # muted
C_TEAL    = RGBColor(0x1A, 0x6B, 0x8A)   # accent teal
C_RED     = RGBColor(0xC0, 0x39, 0x2B)   # heading red
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)   # submitted-by green
C_ORANGE  = RGBColor(0xD2, 0x69, 0x1E)   # orange
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL   = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL_BORDER = RGBColor(0xCC, 0xC4, 0xB8)

W = Inches(13.33)  # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def bg_fill(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_textbox(slide, text, l, t, w, h,
                font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                italic=False, wrap=True, font_name="Inter"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.name    = font_name
    run.font.color.rgb = color if color else C_DARK
    return txb


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0.75)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.color.rgb = line_color if line_color else C_PANEL_BORDER
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_top_stripe(slide):
    """Coloured bar at very top"""
    bar = slide.shapes.add_shape(1, 0, 0, W, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TEAL
    bar.line.fill.background()


def add_slide_label(slide, text, num, total):
    add_textbox(slide, f"{text}   |   {num}/{total}",
                Inches(.2), H - Inches(.35), Inches(4), Inches(.3),
                font_size=9, color=C_MUTED, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════
# SLIDE 1: TITLE (Academic Cover)
# ══════════════════════════════════════════════════════════
s1 = blank_slide(prs)
bg_fill(s1)
add_top_stripe(s1)

# White card background
ml, mt = Inches(1.5), Inches(0.55)
mw, mh = Inches(10.33), Inches(6.5)
add_rect(s1, ml, mt, mw, mh, fill_color=C_WHITE, line_color=C_PANEL_BORDER)

cy = mt + Inches(0.3)

# Project title
add_textbox(s1, "DUAL MODE VISION-GUIDED ROBOTIC ARM",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.55),
            font_size=20, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
cy += Inches(0.58)

add_textbox(s1, "A Project report submitted in partial fulfilment of the requirement for the degree of the",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.35),
            font_size=10, italic=True, color=C_MID, align=PP_ALIGN.CENTER)
cy += Inches(0.37)

add_textbox(s1, "BACHELOR OF TECHNOLOGY",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.32),
            font_size=13, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
cy += Inches(0.30)

add_textbox(s1, "IN",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.28),
            font_size=11, color=C_MID, align=PP_ALIGN.CENTER)
cy += Inches(0.27)

add_textbox(s1, "ELECTRONICS AND COMMUNICATION ENGINEERING",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.32),
            font_size=13, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
cy += Inches(0.35)

add_textbox(s1, "SUBMITTED BY",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.28),
            font_size=10, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
cy += Inches(0.30)

# Team table
team = [
    ("CHADARAM MAANAS",            "Reg No: 321106512004"),
    ("BORA ABHINAV",               "Reg No: 321106512005"),
    ("CHEPURI NATRAJ",             "Reg No: 321106512006"),
    ("MANNAVA SREE PARDHEEV KUMAR","Reg No: 321106512028"),
]
for name, reg in team:
    add_textbox(s1, name,
                ml+Inches(1.2), cy, Inches(4.5), Inches(0.25),
                font_size=11, bold=False, color=C_DARK, align=PP_ALIGN.LEFT)
    add_textbox(s1, reg,
                ml+Inches(5.8), cy, Inches(3.5), Inches(0.25),
                font_size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)
    cy += Inches(0.27)

cy += Inches(0.08)

add_textbox(s1, "UNDER THE ESTEEMED GUIDANCE OF",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.28),
            font_size=10, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
cy += Inches(0.28)

add_textbox(s1, "Prof. P. RAJESH KUMAR M.E., Ph.D.",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.30),
            font_size=13, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
cy += Inches(0.30)

add_textbox(s1, "Professor",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.25),
            font_size=10, color=C_MID, align=PP_ALIGN.CENTER)
cy += Inches(0.3)

# AU Logo
logo_path = r"c:\Users\NATRAJ\Desktop\PPT\Media\AULogo.png"
slide_w_center = W / 2
logo_h = Inches(0.8)
logo_w = Inches(0.8) # Approx aspect ratio 1:1, will adjust based on actual size.
s1.shapes.add_picture(logo_path, slide_w_center - (logo_w/2), cy, height=logo_h)
cy += logo_h + Inches(0.15)

# Divider line
add_rect(s1, ml+Inches(2.5), cy, Inches(5.33), Pt(1),
         fill_color=C_PANEL_BORDER, line_color=C_PANEL_BORDER, line_width=Pt(0))
cy += Inches(0.12)

add_textbox(s1, "DEPARTMENT OF ELECTRONICS AND COMMUNICATION ENGINEERING",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.28),
            font_size=10, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
cy += Inches(0.27)

add_textbox(s1, "ANDHRA UNIVERSITY COLLEGE OF ENGINEERING",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.27),
            font_size=10, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
cy += Inches(0.26)

add_textbox(s1, "VISAKHAPATNAM – 530003",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.25),
            font_size=10, color=C_MID, align=PP_ALIGN.CENTER)
cy += Inches(0.24)

add_textbox(s1, "2022–2026",
            ml+Inches(.3), cy, mw-Inches(.6), Inches(.25),
            font_size=10, color=C_MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# HELPER: generic content slide builder
# ══════════════════════════════════════════════════════════
def title_slide_header(slide, title, subtitle="", slide_n=2, total=10):
    add_top_stripe(slide)
    add_textbox(slide, title,
                Inches(.6), Inches(.35), Inches(12.1), Inches(.65),
                font_size=26, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(.6), Inches(.95), Inches(12.1), Inches(.35),
                    font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_slide_label(slide, title, slide_n, total)


def add_card(slide, l, t, w, h):
    return add_rect(slide, l, t, w, h, fill_color=C_WHITE, line_color=C_PANEL_BORDER)


# ══════════════════════════════════════════════════════════
# SLIDE 2: Abstract
# ══════════════════════════════════════════════════════════
s2 = blank_slide(prs)
bg_fill(s2)
title_slide_header(s2, "Abstract", slide_n=2)

panels = [
    ("🎯 OBJECTIVE",
     "This project develops a 6-DOF robotic arm capable of autonomous object sorting via computer vision and human-in-the-loop gesture control, switching seamlessly between modes without physical intervention."),
    ("🏗️ ARCHITECTURE",
     "Hardware: 6× MG996R servos on a metal arm kit, ESP8266 microcontroller, MPU6050 IMU. Software: ROS2 Jazzy, Gazebo Harmonic for physics simulation, RViz for visualisation, YOLOv8 for hand/object detection."),
    ("📊 PROGRESS",
     "Milestones complete: (1) hardware procured; (2) custom URDF developed after iterative adaptation; (3) arm simulated in Gazebo & RViz; (4) initial ESP8266 servo control written; (5) YOLO hand detection tested — pipeline integration ongoing."),
    ("🔬 SIGNIFICANCE",
     "Dual-track development (hardware + software simulation) validates each track independently, reducing integration risk and ensuring a testable prototype at every phase of the project lifecycle."),
]

positions = [
    (Inches(0.45), Inches(1.35)),
    (Inches(6.85), Inches(1.35)),
    (Inches(0.45), Inches(4.15)),
    (Inches(6.85), Inches(4.15)),
]

for (label, text), (lx, ly) in zip(panels, positions):
    add_card(s2, lx, ly, Inches(6.2), Inches(2.65))
    add_textbox(s2, label, lx+Inches(.2), ly+Inches(.15), Inches(5.8), Inches(.3),
                font_size=9, bold=True, color=C_TEAL)
    add_textbox(s2, text,  lx+Inches(.2), ly+Inches(.45), Inches(5.8), Inches(2.05),
                font_size=11.5, color=C_MID, wrap=True)


# ══════════════════════════════════════════════════════════
# SLIDE 3: System Overview
# ══════════════════════════════════════════════════════════
s3 = blank_slide(prs)
bg_fill(s3)
title_slide_header(s3, "System Architecture & Tools",
                   "Two parallel development tracks converging into one unified system", slide_n=3)

hw_items = ["6-DOF Metal Servo Arm Kit","6× MG996R Servo Motors","ESP8266 Microcontroller",
            "IMU Sensor (MPU6050)","5V / 10A Dedicated PSU"]
sw_items = ["ROS2 Jazzy Jalisco","Gazebo Harmonic (Simulation)","RViz (Visualisation)",
            "YOLOv8 (Hand Detection)","OpenCV + MediaPipe"]

for col_idx, (col_title, items, col_color) in enumerate([
    ("Hardware Track", hw_items, C_ORANGE),
    ("Software Track", sw_items, C_TEAL),
]):
    lx = Inches(0.45) + col_idx * Inches(6.5)
    add_card(s3, lx, Inches(1.45), Inches(6.2), Inches(4.5))
    add_textbox(s3, col_title, lx+Inches(.2), Inches(1.6), Inches(5.8), Inches(.35),
                font_size=14, bold=True, color=col_color)
    for i, item in enumerate(items):
        add_textbox(s3, f"• {item}",
                    lx+Inches(.25), Inches(2.05)+Inches(i*0.58), Inches(5.7), Inches(.5),
                    font_size=12, color=C_MID)

flow = "IMU / Camera Feed  →  ESP8266 / ROS2  →  Servo Motor Commands  →  Arm Motion"
add_textbox(s3, flow, Inches(1.2), Inches(6.15), Inches(11), Inches(.37),
            font_size=11, color=C_TEAL, align=PP_ALIGN.CENTER, bold=True)


# ══════════════════════════════════════════════════════════
# SLIDES 4–8: Process slides (one per phase)
# ══════════════════════════════════════════════════════════
process_slides_data_1 = [
    {
        "num": 4, "title": "Phase 1 — Hardware Procurement", "status_label": "✅ Completed",
        "image": r"c:\Users\NATRAJ\Desktop\PPT\Media\UNASSEMBLED ROBOTIC MODEL.jpeg",
        "intro": "All primary hardware components were ordered and received. Component selection was finalised based on torque, availability, and ESP8266 compatibility.",
        "outcome": "All hardware is on hand and ready for wiring and firmware integration.",
        "steps": [
            ("01", "Component Selection",
             "Six MG996R servos (9–11 kg·cm), ESP8266 NodeMCU for Wi-Fi PWM control, MPU6050 IMU for 6-axis feedback. Initial plan called for 4 servos; upgraded to 6 for full DOF."),
            ("02", "Procurement & Delivery",
             "Components sourced and ordered. Arm kit, servos, ESP8266, IMU, and 5V/10A PSU all received and inspected."),
            ("03", "Inventory Verified",
             "6× MG996R, 1× ESP8266 NodeMCU, 1× MPU6050, 1× 6-DOF frame + claw gripper, 1× 5V/10A PSU — confirmed complete."),
        ]
    }
]

fullimage_slides = [
    (r"c:\Users\NATRAJ\Desktop\PPT\Media\HARWARE REPRESENTAION.jpeg", 5),
    (r"c:\Users\NATRAJ\Desktop\PPT\Media\ASSMEBLED.jpg", 6)
]

process_slides_data_2 = [
    {
        "num": 7, "title": "Phase 2 — 3D Model & URDF Development", "status_label": "✅ Completed (Multiple Iterations)",
        "image": None,
        "intro": "A URDF is mandatory for Gazebo simulation. No pre-existing URDF matched our hardware, so one was created from an adapted open-source reference.",
        "outcome": "Working Xacro URDF produced and integrated into the ROS2 workspace package my_arm_description.",
        "steps": [
            ("01", "Challenge: No Matching URDF",
             "The purchased 6-DOF metal arm kit has no freely available 3D model or URDF online. Without a URDF, simulation and motion planning are impossible."),
            ("02", "Reference Model Located",
             "ferasboulala/five-dof-robot-arm on GitHub was used as a base. Mesh proportions and joint placements were studied against the physical hardware."),
            ("03", "Iterative Adaptation",
             "Multiple attempts: adjusting link lengths, correcting joint axes, converting to ROS2 Xacro format, adding the 6th joint and gripper, fixing mesh collision/inertia errors."),
            ("04", "Working Model Achieved",
             "Stable Xacro URDF produced with 6 correctly defined joints, link inertias, and collision geometry — sufficient for simulation and motion planning."),
        ]
    },
    {
        "num": 8, "title": "Phase 3 — Simulation in Gazebo & RViz", "status_label": "✅ Completed",
        "image": r"c:\Users\NATRAJ\Desktop\PPT\Media\RVIZ IMAGE OF MODEL.png",
        "intro": "With the URDF ready, the arm was loaded into Gazebo Harmonic (physics sim) and RViz (kinematic visualisation) as the primary testbed before hardware deployment.",
        "outcome": "Arm simulates correctly in Gazebo with physics-accurate joint motion. Trajectory commands execute reliably via ROS2 topics.",
        "steps": [
            ("01", "RViz Visualisation",
             "Loaded via display.launch.py. joint_state_publisher_gui sliders confirmed all 6 joints move correctly. Link transforms verified."),
            ("02", "Gazebo Harmonic Launch",
             "Arm spawned via gazebo.launch.py. gz-ros2-bridge configured to relay joint states. JointTrajectoryController from ros2_control activated."),
            ("03", "Controller Verification",
             "ros2 control list_controllers confirmed joint_trajectory_controller and joint_state_broadcaster active."),
            ("04", "Motion Testing",
             "Test trajectory commands published via ros2 topic pub. Pick-place sequences and joint-level position commands prototyped in simulation."),
        ]
    },
    {
        "num": 9, "title": "Phase 4 — Servo Motor Control (ESP8266)", "status_label": "🔄 In Progress",
        "image": None,
        "intro": "Initial development of servo motor control firmware on the ESP8266 began in parallel with simulation. This is the hardware counterpart to the simulated JointTrajectoryController.",
        "outcome": "Individual servo control functional. Multi-servo coordination and ROS2–ESP8266 bridge are the next steps.",
        "steps": [
            ("01", "ESP8266 PWM Control",
             "ESP8266 NodeMCU selected for built-in Wi-Fi and PWM output. Will receive joint angle commands and translate them to PWM signals for all 6 MG996R servos."),
            ("02", "Signal Generation",
             "MG996R servos need 50 Hz PWM, 500–2500 µs pulse width for full motion range. Initial firmware maps angle → pulse width and generates signals on digital IO pins."),
            ("03", "Initial Tests",
             "Preliminary tests on individual servo motors complete. Smooth angle transitions with configurable speed implemented to prevent jerky motion."),
        ]
    },
    {
        "num": 10, "title": "Phase 5 — Computer Vision Integration (YOLO)", "status_label": "🔄 Under Development",
        "image": r"c:\Users\NATRAJ\Desktop\PPT\Media\YOLO HAND DETECTION.jpg",
        "intro": "A YOLOv8-based pipeline was developed to detect a human hand and convert that detection into servo movement commands. This is the core intelligence of gesture control mode.",
        "outcome": "Hand detection functional in isolation. Full pipeline integration into ROS2 + servo chain is the current blocker under active debugging.",
        "steps": [
            ("01", "YOLOv8 Model Setup",
             "YOLOv8 (Ultralytics) configured on NVIDIA RTX 3060 with CUDA 12.1 for accelerated inference. Targeting real-time >30 FPS performance."),
            ("02", "Detection → Command Pipeline",
             "Intended: camera → YOLOv8 bounding box → centroid → IK → 6 joint angles → ESP8266 servo commands. Partially implemented."),
            ("03", "Integration Issues Found",
             "Latency spikes in detection-to-command loop, coordinate mapping calibration errors, and serialization issues passing YOLO output into ROS2 topics."),
            ("04", "Current Status",
             "YOLO detects hands reliably in isolation. End-to-end integration from detection to servo command under active debugging and expected to resolve in next cycle."),
        ]
    },
]


def build_process_slide(pd):
    ps = blank_slide(prs)
    bg_fill(ps)
    title_slide_header(ps, pd["title"], slide_n=pd["num"], total=13)

    # Status badge
    s_color = C_TEAL if "Completed" in pd["status_label"] else C_ORANGE
    add_textbox(ps, pd["status_label"], Inches(5.5), Inches(.9), Inches(3), Inches(.28),
                font_size=9, bold=True, color=s_color, align=PP_ALIGN.CENTER)

    # Left: intro + outcome
    has_img = pd.get("image") is not None
    left_w = Inches(3.0) if has_img else Inches(3.4)
    add_card(ps, Inches(0.4), Inches(1.35), left_w, Inches(2.2))
    add_textbox(ps, "OVERVIEW", Inches(.55), Inches(1.45), left_w - Inches(.3), Inches(.28),
                font_size=8, bold=True, color=C_TEAL)
    add_textbox(ps, pd["intro"], Inches(.55), Inches(1.72), left_w - Inches(.3), Inches(1.65),
                font_size=10.5, color=C_MID, wrap=True)

    add_textbox(ps, "OUTCOME", Inches(.55), Inches(3.65), left_w - Inches(.3), Inches(.25),
                font_size=8, bold=True, color=C_GREEN)
    add_textbox(ps, pd["outcome"], Inches(.55), Inches(3.9), left_w - Inches(.3), Inches(1.2),
                font_size=10, color=RGBColor(0x2D,0x6A,0x4F), wrap=True)

    # Right or Center: steps
    steps_left = Inches(3.6) if has_img else Inches(4.0)
    steps_width = Inches(5.5) if has_img else Inches(9.0)
    
    step_h = min(Inches(1.3), Inches(5.5 / len(pd["steps"])))
    for si, (num, heading, detail) in enumerate(pd["steps"]):
        sy = Inches(1.35) + si * (step_h + Inches(0.1))
        add_card(ps, steps_left, sy, steps_width, step_h)
        add_textbox(ps, num,     steps_left + Inches(.15), sy+Inches(.12), Inches(.5),  step_h-Inches(.12),
                    font_size=9, bold=True, color=C_TEAL)
        add_textbox(ps, heading, steps_left + Inches(.7),  sy+Inches(.1),  steps_width - Inches(.85), Inches(.28),
                    font_size=12, bold=True, color=C_DARK)
        add_textbox(ps, detail,  steps_left + Inches(.7),  sy+Inches(.38), steps_width - Inches(.85), step_h-Inches(.42),
                    font_size=10, color=C_MUTED, wrap=True)

    # Right (if img exists): Image
    if has_img:
        ps.shapes.add_picture(pd["image"], Inches(9.3), Inches(1.35), width=Inches(3.8))

for pd in process_slides_data_1:
    build_process_slide(pd)

# Full Image Slides
for img_path, s_num in fullimage_slides:
    fs = blank_slide(prs)
    bg_fill(fs)
    img_shape = fs.shapes.add_picture(img_path, 0, 0, height=H)
    # Center image horizontally
    img_shape.left = int((W - img_shape.width) / 2)


for pd in process_slides_data_2:
    build_process_slide(pd)


# ══════════════════════════════════════════════════════════
# SLIDE 9: Challenges & Status
# ══════════════════════════════════════════════════════════
s11 = blank_slide(prs)
bg_fill(s11)
title_slide_header(s11, "Challenges & Current Status",
                   "Dual-track development: each track validated independently, converging at integration", slide_n=11, total=12)

hw_done    = ["All components procured","Specs verified","ESP8266 PWM firmware written","Individual servo response tested"]
hw_pending = ["Full arm assembly & wiring","Multi-servo coordinated control","ROS2 ↔ ESP8266 serial bridge","IMU integration"]
sw_done    = ["URDF created after iterative adaptation","Arm visualised in RViz","Gazebo simulation running",
              "ros2_control JTC active","YOLOv8 detection working in isolation"]
sw_pending = ["YOLO → ROS2 command pipeline","Camera-to-workspace calibration","MoveIt2 motion planning","Autonomous mode logic"]

for col_idx, (col_title, done, pending, col_color) in enumerate([
    ("Hardware Track", hw_done, hw_pending, C_ORANGE),
    ("Software Track", sw_done, sw_pending, C_TEAL),
]):
    lx = Inches(0.4) + col_idx * Inches(6.55)
    add_card(s11, lx, Inches(1.4), Inches(6.25), Inches(4.65))
    add_textbox(s11, col_title, lx+Inches(.18), Inches(1.5), Inches(5.9), Inches(.35),
                font_size=14, bold=True, color=col_color)
    add_textbox(s11, "COMPLETED", lx+Inches(.18), Inches(1.9), Inches(5.9), Inches(.25),
                font_size=9, bold=True, color=C_MUTED)
    cur_y = Inches(2.15)
    for item in done:
        add_textbox(s11, f"✅  {item}", lx+Inches(.18), cur_y, Inches(5.9), Inches(.3),
                    font_size=11, color=C_MID)
        cur_y += Inches(0.31)
    add_textbox(s11, "PENDING", lx+Inches(.18), cur_y+Inches(.05), Inches(5.9), Inches(.25),
                font_size=9, bold=True, color=C_MUTED)
    cur_y += Inches(0.3)
    for item in pending:
        add_textbox(s11, f"⏳  {item}", lx+Inches(.18), cur_y, Inches(5.9), Inches(.3),
                    font_size=11, color=C_MUTED)
        cur_y += Inches(0.31)

key_text = ("Key Challenge: The most significant obstacle was the absence of a ready-made URDF "
            "for the purchased arm kit, requiring multiple iterative reconstruction attempts before "
            "a working simulation model was produced.")
add_textbox(s11, key_text, Inches(0.4), Inches(6.2), Inches(12.6), Inches(.45),
            font_size=11, color=C_ORANGE, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 10: Next Steps / Roadmap
# ══════════════════════════════════════════════════════════
s12 = blank_slide(prs)
bg_fill(s12)
title_slide_header(s12, "Next Steps & Roadmap",
                   "Structured development path — testable artefact at every phase", slide_n=12, total=12)

phases = [
    ("P1", "Hardware Assembly",       C_GREEN,  ["Wire all 6 servos","Connect IMU","PSU integration","Bench test all joints"]),
    ("P2", "ROS2 ↔ Hardware Bridge",  C_TEAL,   ["ESP8266 serial bridge node","ROS2 JointTrajectory → servo","Real HW controller test","ros2_control HW interface"]),
    ("P3", "Vision Pipeline Fix",     C_RED,    ["Debug YOLO → ROS2 pipeline","Camera calibration","Coordinate mapping","Latency test"]),
    ("P4", "Full System Integration", C_ORANGE, ["Autonomous sorting mode","Gesture control (local)","Mode switching logic","50-cycle stress test"]),
]

for col_idx, (pid, plabel, pcol, tasks) in enumerate(phases):
    lx = Inches(0.35) + col_idx * Inches(3.28)
    add_card(s12, lx, Inches(1.45), Inches(3.1), Inches(4.8))
    add_textbox(s12, pid,    lx+Inches(.15), Inches(1.55), Inches(2.8), Inches(.25),
                font_size=9, color=C_MUTED, bold=True)
    add_textbox(s12, plabel, lx+Inches(.15), Inches(1.78), Inches(2.8), Inches(.38),
                font_size=13, bold=True, color=pcol)
    for ti, task in enumerate(tasks):
        add_textbox(s12, f"→  {task}",
                    lx+Inches(.15), Inches(2.2)+Inches(ti*.52), Inches(2.8), Inches(.45),
                    font_size=11, color=C_MID)

note = "Each phase produces a standalone demonstrable prototype — ensuring progress is always presentable regardless of final timeline."
add_textbox(s12, note, Inches(1.2), Inches(6.5), Inches(11), Inches(.38),
            font_size=11, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
out_path = r"c:\Users\NATRAJ\Desktop\PPT\presentation.pptx"
prs.save(out_path)
print(f"✅  Saved → {out_path}")
